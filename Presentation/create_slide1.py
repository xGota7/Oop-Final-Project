"""Generate Quiz Arena architecture overview slide (editable PowerPoint)."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE

OUT = Path(__file__).resolve().parent / "QuizArena.pptx"

NAVY = RGBColor(3, 40, 77)
BLUE = RGBColor(0, 76, 158)
LIGHT_BLUE = RGBColor(236, 246, 255)
GREEN = RGBColor(37, 106, 37)
LIGHT_GREEN = RGBColor(240, 250, 236)
PURPLE = RGBColor(95, 47, 150)
LIGHT_PURPLE = RGBColor(246, 238, 255)
ORANGE = RGBColor(205, 64, 15)
LIGHT_ORANGE = RGBColor(255, 242, 231)
TEAL = RGBColor(0, 126, 128)
LIGHT_TEAL = RGBColor(232, 250, 250)
GRAY = RGBColor(80, 90, 110)
LIGHT_GRAY = RGBColor(245, 248, 252)
MID_GRAY = RGBColor(195, 205, 218)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

BODY_FONT = "Calibri"
ICON_FONT = "Segoe UI Symbol"
CODE_FONT = "Consolas"
MIN_SIZE = 0.05

G_CAP = "\u2691"
G_LIST = "\u2263"
G_FILE = "\u25a4"
G_GAME = "\u25c8"
G_Q = "?"
G_SCREEN = "\u25ad"
G_DISK = "\u25a3"
G_PERSON = "\u25cf"
G_CUP = "\u265c"


def _size(v):
    return Inches(max(float(v), MIN_SIZE))


def set_text_font(run, size=16, color=NAVY, bold=False, italic=False, font=None):
    run.font.name = font or BODY_FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_text(slide, text, x, y, w, h, size=16, color=NAVY, bold=False,
             italic=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), _size(w), _size(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    for i, chunk in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = chunk
        set_text_font(r, size, color, bold, italic, font)
    return shape


def add_icon(slide, glyph, x, y, w, h, size=22, color=NAVY, align=PP_ALIGN.CENTER):
    return add_text(slide, glyph, x, y, w, h, size, color, True,
                    align=align, valign=MSO_ANCHOR.MIDDLE, font=ICON_FONT)


def add_round_rect(slide, x, y, w, h, fill=WHITE, line=NAVY, radius=True, width=1.25):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), _size(w), _size(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if width <= 0:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(width)
    if radius:
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass
    return shp


def add_plain_rect(slide, x, y, w, h, fill=WHITE, line=NAVY, width=1.0):
    return add_round_rect(slide, x, y, w, h, fill, line, radius=False, width=width)


def add_line(slide, x1, y1, x2, y2, color=NAVY, width=1.5, dashed=False, arrow=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if dashed:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if arrow:
        try:
            line.line.end_arrowhead = True
        except Exception:
            pass
    return line


def add_elbow(slide, x1, y1, x2, y2, color=NAVY, width=1.5, arrow=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.ELBOW, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if arrow:
        try:
            line.line.end_arrowhead = True
        except Exception:
            pass
    return line


def add_diamond(slide, cx, cy, size=0.16, color=NAVY):
    d = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND, Inches(cx - size / 2), Inches(cy - size / 2),
        Inches(size), Inches(size))
    d.fill.solid()
    d.fill.fore_color.rgb = color
    d.line.color.rgb = color
    d.line.width = Pt(0.75)
    return d


def add_icon_circle(slide, glyph, x, y, d=0.42, fill=NAVY, color=WHITE, size=17):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.color.rgb = fill
    add_icon(slide, glyph, x, y + 0.01, d, d - 0.01, size, color)
    return c


def add_class_box(slide, title, lines, x, y, w, h, fill, line, icon=None,
                  title_size=18, body_size=12):
    add_round_rect(slide, x, y, w, h, fill, line, True, 1.35)
    header = min(0.52, h * 0.60)
    tx = x + 0.22
    if icon:
        add_icon(slide, icon, x + 0.20, y + 0.06, 0.50, header - 0.06, 20, line)
        tx = x + 0.76
    add_text(slide, title, tx, y + 0.06, w - (tx - x) - 0.16, header - 0.04,
             title_size, line, True, valign=MSO_ANCHOR.MIDDLE)
    if lines:
        add_line(slide, x + 0.18, y + header + 0.06, x + w - 0.18, y + header + 0.06, line, 1.0)
        add_text(slide, lines, x + 0.22, y + header + 0.14, w - 0.44,
                 h - header - 0.20, body_size, BLACK, font=CODE_FONT)


def add_footer(slide, num):
    add_plain_rect(slide, 0, 8.25, 16, 0.75, NAVY, NAVY, 0)
    add_text(slide, str(num), 14.90, 8.36, 0.80, 0.40, 21, WHITE, True,
             align=PP_ALIGN.RIGHT)


def build():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Soft background
    add_plain_rect(slide, 0, 0, 16, 9, LIGHT_GRAY, LIGHT_GRAY, 0)

    # Title
    add_round_rect(slide, 0.22, 0.16, 0.10, 1.25, NAVY, NAVY, True, 0)
    add_text(slide, "Quiz Arena", 0.55, 0.18, 11.5, 0.58, 42, NAVY, True)
    add_text(slide, "C++ Object-Oriented Programming Final Project",
             0.57, 0.82, 11.0, 0.32, 16, RGBColor(52, 68, 95), italic=True)
    add_text(slide,
             "Text-based quiz game with saves, leaderboard, and polymorphic question types.",
             0.58, 1.18, 11.0, 0.28, 13, NAVY)

    # Badge
    add_round_rect(slide, 12.55, 0.30, 2.95, 0.78, LIGHT_BLUE, MID_GRAY, True, 1.0)
    add_icon(slide, G_CAP, 12.72, 0.48, 0.36, 0.30, 18, NAVY)
    add_text(slide, "Prepared for", 13.16, 0.40, 2.1, 0.20, 10, GRAY)
    add_text(slide, "OOP Final Project", 13.16, 0.62, 2.1, 0.28, 13, NAVY, True)

    # Agenda
    add_round_rect(slide, 0.45, 1.70, 4.50, 5.85, LIGHT_BLUE, MID_GRAY, True, 1.0)
    add_icon(slide, G_LIST, 0.75, 1.95, 0.45, 0.40, 22, NAVY)
    add_text(slide, "Agenda", 1.35, 1.95, 3.2, 0.42, 24, NAVY, True)
    add_line(slide, 0.75, 2.55, 4.60, 2.55, MID_GRAY, 1.0)

    agenda = [
        "Project idea and rules",
        "Question hierarchy",
        "Player and runtime state",
        "Save system",
        "Leaderboard",
        "How everything connects in QuizGame",
    ]
    yy = 2.85
    for i, item in enumerate(agenda, 1):
        add_icon_circle(slide, str(i), 0.80, yy - 0.05, 0.40, NAVY, WHITE, 14)
        add_text(slide, item, 1.40, yy - 0.02, 3.25, 0.55, 14, NAVY, valign=MSO_ANCHOR.MIDDLE)
        if i < 6:
            add_line(slide, 1.30, yy + 0.55, 4.55, yy + 0.55, MID_GRAY, 0.8, dashed=True)
        yy += 0.70

    # Connectors (drawn under boxes)
    add_line(slide, 9.10, 2.55, 9.10, 3.05, BLACK, 1.7, arrow=True)
    add_line(slide, 11.00, 3.70, 12.55, 3.55, BLACK, 1.3, dashed=True, arrow=True)
    add_line(slide, 11.00, 4.30, 12.55, 4.65, BLACK, 1.3, dashed=True, arrow=True)
    add_line(slide, 11.00, 4.90, 12.55, 6.10, BLACK, 1.3, dashed=True, arrow=True)
    add_elbow(slide, 9.10, 5.35, 7.50, 6.70, BLACK, 1.4, arrow=True)
    add_elbow(slide, 9.10, 5.35, 10.70, 6.70, BLACK, 1.4, arrow=True)
    add_diamond(slide, 9.10, 5.42)

    # main.cpp
    add_class_box(slide, "main.cpp", "int main()", 8.05, 1.62, 2.10, 0.94,
                  LIGHT_PURPLE, PURPLE, icon=G_FILE, title_size=14, body_size=11)

    # QuizGame
    add_class_box(
        slide, "QuizGame",
        "run()\nvector<Question*> m_questions\nPlayer m_player\n"
        "Leaderboard m_leaderboard\nm_difficulty / m_targetScore",
        7.20, 3.05, 3.80, 2.30, LIGHT_BLUE, BLUE, icon=G_GAME,
        title_size=20, body_size=11)

    add_text(slide, "uses", 11.40, 3.40, 0.70, 0.22, 10, BLACK, italic=True)
    add_text(slide, "uses", 11.40, 4.20, 0.70, 0.22, 10, BLACK, italic=True)
    add_text(slide, "uses", 11.10, 5.15, 0.70, 0.22, 10, BLACK, italic=True)

    add_class_box(
        slide, "Question (Hierarchy)",
        "abstract base class\nMultipleChoiceQuestion /\nTrueFalseQuestion",
        12.55, 2.95, 3.00, 1.30, LIGHT_GREEN, GREEN, icon=G_Q,
        title_size=14, body_size=10)

    add_class_box(slide, "ConsoleUI", "all input / output",
                  12.55, 4.40, 3.00, 0.95, LIGHT_PURPLE, PURPLE, icon=G_SCREEN,
                  title_size=15, body_size=11)

    add_class_box(
        slide, "SaveManager",
        "save / load game data\n(slots, progress, state)",
        12.55, 5.55, 3.00, 1.20, LIGHT_ORANGE, ORANGE, icon=G_DISK,
        title_size=15, body_size=11)

    # Composition targets (accurate to QuizGame.h — no duplicate Player)
    add_class_box(slide, "Player", "m_name, m_score, m_lives\ncurrent state",
                  6.35, 6.70, 2.30, 1.15, LIGHT_TEAL, TEAL, icon=G_PERSON,
                  title_size=15, body_size=10)
    add_text(slide, "1", 7.55, 6.42, 0.25, 0.22, 12, BLACK)

    add_class_box(slide, "Leaderboard",
                  "unordered_map scores\ntop scores by mode",
                  9.55, 6.70, 2.30, 1.15, LIGHT_BLUE, BLUE, icon=G_CUP,
                  title_size=15, body_size=10)
    add_text(slide, "1", 10.60, 6.42, 0.25, 0.22, 12, BLACK)

    add_footer(slide, 1)
    prs.save(str(OUT))
    print("Created:", OUT)


if __name__ == "__main__":
    build()
